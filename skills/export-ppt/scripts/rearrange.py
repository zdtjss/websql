#!/usr/bin/env python3
"""
PPT Template Rearrange — 基于模板生成新的演示文稿
从模板中按索引选取并重排幻灯片，支持重复和去重。
基于 Anthropic pptx skill 的 rearrange.py 思想。
"""

import sys
import os
import copy
from lxml import etree
from pptx import Presentation
from pptx.parts.slide import SlidePart


_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.dirname(_SCRIPT_DIR))


NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def rearrange(template_path, output_path, slide_indices):
    """从模板中选取指定索引的幻灯片重新排列为新演示文稿"""
    template = Presentation(template_path)
    total = len(template.slides)

    for idx in slide_indices:
        if idx < 0 or idx >= total:
            print(f"ERROR: 幻灯片索引 {idx} 超出范围 (0-{total - 1})", file=sys.stderr)
            sys.exit(1)

    output = Presentation(template_path)

    xml_slides = output.part.presentation_element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
    )
    if xml_slides is None:
        root = output.part.presentation_element
        xml_slides = etree.SubElement(root, "{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst")

    for child in list(xml_slides):
        xml_slides.remove(child)

    slide_id = 256
    rId_counter = 1

    for template_idx in slide_indices:
        src_slide = template.slides[template_idx]
        src_part = src_slide.part

        new_part = copy.deepcopy(src_part)

        rId = f"rId{rId_counter}"
        output.part.relate_to(new_part, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")

        sld_elem = etree.SubElement(xml_slides, "{http://schemas.openxmlformats.org/presentationml/2006/main}sldId")
        sld_elem.set("id", str(slide_id))
        sld_elem.set("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id", rId)

        slide_id += 1
        rId_counter += 1

    output.save(output_path)
    print(f"已生成 {len(slide_indices)} 张幻灯片 -> {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("用法: python rearrange.py <template.pptx> <output.pptx> <idx1,idx2,...>", file=sys.stderr)
        print("示例: python rearrange.py template.pptx output.pptx 0,3,3,5", file=sys.stderr)
        sys.exit(1)

    template_path = sys.argv[1]
    output_path = sys.argv[2]
    indices = [int(x.strip()) for x in sys.argv[3].split(",")]
    rearrange(template_path, output_path, indices)
