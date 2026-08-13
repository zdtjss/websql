# -*- coding: utf-8 -*-
"""export_ppt.py — 模板驱动的 PPT 导出脚本（渲染层唯一引擎）。

版式定义在 skills/export-ppt/templates/slides_template.pptx（母版 + 版式），
本脚本只做内容填充与布局自适应，不再手工坐标排版。
换肤/换版式：直接编辑模板文件（或改 skills/lib/build_templates.py 后重新生成）。

版式索引约定（与 build_templates.py 对齐）：
    L_COVER=0 封面   L_CONTENT=1 内容/目录/表格/KPI   L_SECTION=2 章节分隔
    L_CHART=5 图表   L_BLANK=6 致谢

JSON 契约（与 SKILL.md 一致）：
    data 模式:    {mode, title, columns, data, highlights, chartPaths, outputPath}
    content 模式: {mode, title, sections:[{title, level, blocks:[{kind|type, ...}]}],
                   markdown?, outputPath}
"""

import json
import math
import os
import sys
import tempfile
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# lib 目录加入 sys.path（兼容从任意工作目录执行）
sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "lib"))

from chart_common import DarkChartStyle, create_chart  # noqa: E402
from md_blocks import markdown_to_sections, normalize_sections  # noqa: E402


class Theme:
    # PPT 颜色（与模板文件配色保持一致）
    BG_DARK = RGBColor(0x0A, 0x16, 0x28)       # 深空蓝背景
    BG_CARD = RGBColor(0x12, 0x23, 0x3D)       # 卡片背景
    BLUE = RGBColor(0x00, 0xA8, 0xFF)          # 科技蓝
    CYAN = RGBColor(0x00, 0xF5, 0xD4)          # 电光青
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0xB0, 0xC4, 0xD8)          # 浅灰蓝
    GOLD = RGBColor(0xFF, 0xD7, 0x00)          # 星光金
    GREEN = RGBColor(0x00, 0xE6, 0x96)         # 增长绿
    RED = RGBColor(0xE6, 0x39, 0x46)           # 警示红

    CHART_BG = '#0A1628'
    CHART_FACE = '#0F1E36'
    CHART_GRID = '#1E3A5F'


_CHART_STYLE = DarkChartStyle()

# 图表页图片区域（与模板 layout[5] 的 chart_frame 几何一致）
CHART_AREA = (1.0, 1.5, 11.3, 5.6)


def _set_ea(rPr, typeface):
    """在 a:latin 之后追加/更新 a:ea 东亚字体。"""
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    ea.set('typeface', typeface)


# ═══════════════════════════════════════════════════════════════
# PPTBuilder：模板驱动渲染
# ═══════════════════════════════════════════════════════════════
class PPTBuilder:
    L_COVER = 0
    L_CONTENT = 1
    L_SECTION = 2
    L_CHART = 5
    L_BLANK = 6

    def __init__(self, template_path=None):
        template_path = template_path or os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            "..", "templates", "slides_template.pptx")
        self.prs = Presentation(template_path)
        self._temp_files = []
        self._setup_matplotlib()

    def _setup_matplotlib(self):
        plt.style.use('dark_background')
        plt.rcParams.update({
            'figure.facecolor': Theme.CHART_FACE,
            'axes.facecolor': Theme.CHART_BG,
            'text.color': '#FFFFFF',
            'axes.labelcolor': '#B0C4D8',
            'xtick.color': '#B0C4D8',
            'ytick.color': '#B0C4D8',
            'axes.edgecolor': Theme.CHART_GRID,
            'grid.color': Theme.CHART_GRID,
            'grid.alpha': 0.3,
        })

    # ─── 基础元素 ───────────────────────────────────────────────

    def _apply_font(self, run, size, color, bold=False):
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = '微软雅黑'
        _set_ea(run._r.get_or_add_rPr(), '微软雅黑')

    def _add_text(self, slide, text, left, top, width, height,
                  size=16, color=Theme.GRAY, bold=False, align=PP_ALIGN.LEFT):
        """动态文本框（用于版式占位符之外的补充内容）。"""
        box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(str(text).split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(6)
            self._apply_font(p.add_run(), size, color, bold)
        return box

    def _remove_placeholder(self, slide, idx):
        """移除占位符（为动态内容让位，如 body 区放表格/卡片）。"""
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx == idx:
                ph._element.getparent().remove(ph._element)
                return

    def _est_lines(self, text, size_pt, width_in):
        """估算文本换行后的总行数（中文按全角宽估算）。"""
        char_w = size_pt / 72.0 * 1.1
        cpl = max(1, int(width_in / char_w))
        total = 0
        for line in str(text).split('\n'):
            total += max(1, math.ceil(max(len(line), 1) / cpl))
        return total

    def _fit_text(self, slide, text, left, top, width, height,
                  base_size=16, min_size=10, color=Theme.GRAY,
                  bold=False, align=PP_ALIGN.LEFT):
        """填充文本并按可用空间自适应降字号（替代手工坐标防溢出）。"""
        size = base_size
        while size > min_size:
            line_h = (size + 6) / 72.0
            if self._est_lines(text, size, width) * line_h <= height:
                break
            size -= 2
        return self._add_text(slide, text, left, top, width, height,
                              size=size, color=color, bold=bold, align=align)

    def _save_chart(self, fig):
        path = os.path.join(tempfile.gettempdir(),
                            f'ppt_chart_{len(self._temp_files)}.png')
        fig.savefig(path, dpi=150, bbox_inches='tight',
                    facecolor=Theme.CHART_FACE, edgecolor='none')
        plt.close(fig)
        self._temp_files.append(path)
        return path

    # ─── 页面类型（全部基于模板版式填充）───────────────────────

    def add_cover(self, title, subtitle="", date="", author=""):
        """封面页：layout[0] 的标题/副标题占位符。"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_COVER])
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        bottom = f"{date}    {author}".strip()
        if bottom:
            self._add_text(slide, bottom, 1.5, 5.05, 6.0, 0.6,
                           size=14, color=Theme.GRAY)

    def add_toc(self, sections):
        """目录页：layout[1]，编号 + 标题。"""
        if not sections:
            return
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_CONTENT])
        slide.shapes.title.text = "目录"
        lines = []
        for i, sec in enumerate(sections):
            num = f"0{i + 1}" if i < 9 else str(i + 1)
            lines.append(f"{num}  {sec}")
        body = slide.placeholders[1]
        body.text = "\n".join(lines)
        for i, p in enumerate(body.text_frame.paragraphs):
            p.space_after = Pt(10)
            for run in p.runs:
                self._apply_font(run, 18, Theme.WHITE)

    def add_section_divider(self, title, desc=""):
        """章节分隔页：layout[2]。"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_SECTION])
        slide.shapes.title.text = title
        if desc and len(slide.placeholders) > 1:
            slide.placeholders[1].text = desc

    def add_text_page(self, title, content, highlight_indices=None):
        """纯文字要点页：layout[1]。"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_CONTENT])
        slide.shapes.title.text = title
        highlight_indices = highlight_indices or []
        lines = []
        for i, item in enumerate(content):
            lines.append(f"▸ {item}" if i in highlight_indices else f"• {item}")
        text = "\n".join(lines)
        body = slide.placeholders[1]
        body.text = text
        for i, p in enumerate(body.text_frame.paragraphs):
            p.space_after = Pt(8)
            color = Theme.CYAN if i in highlight_indices else Theme.GRAY
            for run in p.runs:
                self._apply_font(run, 16, color)

    def add_kpi_page(self, title, kpis):
        """KPI 卡片页：layout[1] 标题 + body 区动态卡片。

        kpis: list of {label, value, change, trend}
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_CONTENT])
        slide.shapes.title.text = title
        self._remove_placeholder(slide, 1)

        n = max(1, len(kpis))
        gap = 0.3
        card_w = min(2.8, (11.5 - gap * (n - 1)) / n)
        for i, kpi in enumerate(kpis):
            x = 0.9 + i * (card_w + gap)
            y = 1.7
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                Inches(card_w), Inches(4.6))
            card.fill.solid()
            card.fill.fore_color.rgb = Theme.BG_CARD
            card.line.color.rgb = Theme.BLUE
            card.line.width = Pt(1)
            self._add_text(slide, str(kpi.get('label', '')),
                           x + 0.25, y + 0.3, card_w - 0.5, 0.5,
                           size=13, color=Theme.GRAY)
            self._add_text(slide, str(kpi.get('value', '')),
                           x + 0.25, y + 1.1, card_w - 0.5, 1.1,
                           size=28, color=Theme.WHITE, bold=True)
            trend = kpi.get('trend', '')
            if trend == 'up':
                t_color, arrow = Theme.GREEN, "▲"
            elif trend == 'down':
                t_color, arrow = Theme.RED, "▼"
            else:
                t_color, arrow = Theme.GRAY, "─"
            self._add_text(slide, f"{arrow} {kpi.get('change', '')}",
                           x + 0.25, y + 2.7, card_w - 0.5, 0.6,
                           size=14, color=t_color)

    def add_chart_page(self, title, chart_type, data, insights=None):
        """图表页：layout[5]，图片按 CHART_AREA 插入。"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_CHART])
        slide.shapes.title.text = title
        fig = self._create_chart(chart_type, data)
        img_path = self._save_chart(fig)

        if insights:
            area = (CHART_AREA[0], CHART_AREA[1], 7.5, CHART_AREA[3])
            slide.shapes.add_picture(img_path, Inches(*area[:2]),
                                     Inches(area[2]), Inches(area[3]))
            text = '\n'.join(f"▸ {s}" for s in insights)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1.7),
                Inches(4.2), Inches(5.2))
            card.fill.solid()
            card.fill.fore_color.rgb = Theme.BG_CARD
            card.line.color.rgb = Theme.BLUE
            card.line.width = Pt(1)
            self._fit_text(slide, text, 9.0, 2.1, 3.7, 4.4,
                           base_size=14, color=Theme.GRAY)
        else:
            slide.shapes.add_picture(img_path, Inches(CHART_AREA[0]),
                                     Inches(CHART_AREA[1]),
                                     Inches(CHART_AREA[2]),
                                     Inches(CHART_AREA[3]))

    def add_table_page(self, title, headers, rows, total=None):
        """表格页：layout[1] 标题 + body 区表格。"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_CONTENT])
        slide.shapes.title.text = title
        self._remove_placeholder(slide, 1)

        num_rows = len(rows) + 1
        num_cols = max(1, len(headers))
        if len(rows) == 0:
            return
        col_w = min(1.9, 11.5 / num_cols)
        tbl_w = col_w * num_cols
        row_h = 0.38
        max_rows = int(5.0 / row_h) - 1
        display = rows[:max_rows]
        tbl_h = row_h * (len(display) + 1)

        table_shape = slide.shapes.add_table(
            len(display) + 1, num_cols,
            Inches((13.333 - tbl_w) / 2), Inches(1.7),
            Inches(tbl_w), Inches(tbl_h))
        table = table_shape.table

        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(h)[:24]
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    self._apply_font(run, 11, Theme.WHITE, bold=True)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x7E)

        for r_idx, row in enumerate(display):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)[:28]
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    for run in p.runs:
                        self._apply_font(run, 10, Theme.GRAY)
                if r_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = Theme.BG_CARD

        if total is not None and total > len(display):
            self._add_text(slide, f"※ 共 {total} 行，以上仅展示前 {len(display)} 行",
                           0.9, 1.7 + tbl_h + 0.15, 8, 0.5,
                           size=12, color=Theme.GRAY)

    def add_thank_you(self, text="感谢聆听", contact=""):
        """致谢页：layout[6] 空白版式 + 动态文本。"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.L_BLANK])
        self._add_text(slide, text, 2, 2.9, 9, 1.2,
                       size=40, color=Theme.WHITE, bold=True,
                       align=PP_ALIGN.CENTER)
        if contact:
            self._add_text(slide, contact, 2, 4.6, 9, 0.8,
                           size=14, color=Theme.GRAY, align=PP_ALIGN.CENTER)

    # ─── 图表生成（共享 chart_common.py）────────────────────────

    def _create_chart(self, chart_type, data):
        return create_chart(_CHART_STYLE, chart_type, data)

    # ─── 保存 ──────────────────────────────────────────────────

    def save(self, filepath):
        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        try:
            self.prs.save(filepath)
        finally:
            for f in self._temp_files:
                try:
                    os.remove(f)
                except OSError:
                    pass
            self._temp_files.clear()


# ═══════════════════════════════════════════════════════════════
# 渲染编排
# ═══════════════════════════════════════════════════════════════

GROUP_COLUMNS = {"表名", "table_name", "TABLE_NAME", "表名称"}


def render_content(builder, sections):
    """content 模式渲染：目录 → 各节（分隔页/内容页/表格页/图表页）。"""
    toc_titles = [s["title"] for s in sections if s["title"]]
    if len(toc_titles) > 1:
        builder.add_toc(toc_titles)

    for sec in sections:
        sec_title = sec.get("title", "")
        level = sec.get("level", 2)
        if level == 1 and sec_title:
            builder.add_section_divider(sec_title)

        page_title = sec_title
        text_lines = []

        def flush_text():
            nonlocal text_lines
            if text_lines:
                builder.add_text_page(page_title, text_lines)
                text_lines = []

        for b in sec.get("blocks", []):
            kind = b.get("kind", "paragraph")
            if kind == "paragraph":
                if b.get("content"):
                    text_lines.append(b["content"])
            elif kind == "subheading":
                flush_text()
                page_title = b.get("content", "") or page_title
            elif kind == "bullets":
                for it in b.get("items", []):
                    text_lines.append(it)
            elif kind == "code":
                text_lines.append("代码：")
                text_lines.extend(b.get("lines", []))
            elif kind == "table":
                flush_text()
                headers = b.get("headers", [])
                rows = b.get("rows", [])
                if headers:
                    builder.add_table_page(page_title or "数据表", headers, rows)
            elif kind == "image":
                flush_text()
                labels = b.get("labels", [])
                values = b.get("values", [])
                if labels and values:
                    chart_data = {
                        "title": b.get("chart_title", ""),
                        "categories": labels,
                        "series": [{"name": b.get("chart_title", ""),
                                    "values": values}],
                    }
                    builder.add_chart_page(page_title or "图表分析",
                                           b.get("chart_type", "bar"), chart_data)
        flush_text()


def render_data(builder, data):
    """data 模式渲染（与原契约逻辑一致）。"""
    columns = data.get("columns", [])
    rows = data.get("data", [])
    highlights = data.get("highlights", [])
    chart_paths = data.get("chartPaths", [])

    if highlights:
        builder.add_text_page("核心发现", [str(h) for h in highlights],
                               highlight_indices=[0])

    if columns and rows:
        group_col = next((c for c in columns if c in GROUP_COLUMNS), None)

        if group_col:
            from collections import OrderedDict
            groups = OrderedDict()
            for row in rows:
                gk = str(row.get(group_col, ""))
                groups.setdefault(gk, []).append(row)
            other_cols = [c for c in columns if c != group_col]
            for gk, gk_rows in groups.items():
                display_rows = [[str(row.get(c, "")) for c in other_cols]
                                for row in gk_rows]
                builder.add_table_page(
                    f"{group_col}: {gk}（{len(gk_rows)} 条）",
                    other_cols, display_rows[:15], total=len(gk_rows))
        else:
            display_rows = [[str(row.get(c, "")) for c in columns]
                            for row in rows]
            # 数据概览图表（取前 10 行、第二个数值列）
            num_cols = []
            for c in columns[1:]:
                vals = [r for r in rows[:10] if _is_number(str(r.get(c, "")))]
                if vals:
                    num_cols.append(c)
                    break
            if num_cols:
                col = num_cols[0]
                chart_data = {
                    "title": "数据分布",
                    "categories": [str(r[0])[:12] for r in display_rows[:10]],
                    "series": [{"name": col, "values": [
                        _to_float(r.get(col, 0)) for r in rows[:10]]}],
                }
                builder.add_chart_page("数据概览", "bar", chart_data)
            builder.add_table_page("数据明细", columns, display_rows[:20],
                                   total=len(rows))

    if chart_paths:
        for cp in chart_paths[:3]:
            if os.path.exists(cp):
                slide = builder.prs.slides.add_slide(
                    builder.prs.slide_layouts[PPTBuilder.L_CHART])
                slide.shapes.title.text = "图表分析"
                slide.shapes.add_picture(cp, Inches(CHART_AREA[0]),
                                         Inches(CHART_AREA[1]),
                                         Inches(CHART_AREA[2]),
                                         Inches(CHART_AREA[3]))


def _is_number(s):
    t = s.replace(",", "").replace("%", "")
    try:
        float(t)
        return True
    except ValueError:
        return False


def _to_float(v):
    try:
        return float(str(v).replace(",", ""))
    except (TypeError, ValueError):
        return 0.0


# ═══════════════════════════════════════════════════════════════
# 入口
# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    # Windows 下默认 stdin/stdout 编码可能为 GBK（cp936），显式重配置为 UTF-8
    try:
        sys.stdin.reconfigure(encoding='utf-8')
        sys.stdout.reconfigure(encoding='utf-8')
    except (AttributeError, ValueError):
        pass

    try:
        raw = sys.stdin.read()
        # execute 工具不支持直接向子进程 stdin 写入，Agent 可能通过
        # `python xxx.py < input.json` 重定向或位置参数传入 JSON 文件
        if not raw and len(sys.argv) > 1:
            try:
                with open(sys.argv[1], "r", encoding="utf-8-sig") as f:
                    raw = f.read()
            except OSError as e:
                print(json.dumps({"success": False,
                                  "error": f"读取输入文件失败: {e}"},
                                 ensure_ascii=False))
                sys.exit(1)
        if not raw:
            print(json.dumps({"success": False, "error": "stdin 为空"},
                             ensure_ascii=False))
            sys.exit(1)
        data = json.loads(raw)
    except Exception as e:
        print(json.dumps({"success": False, "error": f"解析输入失败: {e}"},
                         ensure_ascii=False))
        sys.exit(1)

    try:
        mode = data.get("mode", "data")
        title = data.get("title", "数据报告")
        output_path = data.get("outputPath", "exports/output.pptx")
        # 规范化输出路径：去掉前导 / 使其成为相对路径，确保 Windows 下
        # 保存到项目 exports/ 目录（/exports/x.pptx 会解析到驱动器根）
        if output_path.startswith("/"):
            output_path = output_path.lstrip("/")

        builder = PPTBuilder()
        builder.add_cover(title, subtitle="数据分析报告",
                          date=datetime.now().strftime("%Y-%m-%d"))

        if mode == "content":
            if data.get("markdown"):
                sections = markdown_to_sections(data["markdown"])
            else:
                sections = normalize_sections(data.get("sections", []))
            render_content(builder, sections)
        else:
            render_data(builder, data)

        builder.add_thank_you()
        builder.save(output_path)

        slide_count = len(builder.prs.slides)
        url_path = "/" + output_path.replace("\\", "/").lstrip("/")
        print(json.dumps({"success": True, "slideCount": slide_count,
                          "outputPath": url_path}, ensure_ascii=False))
    except Exception as e:
        safe_err = str(e).encode('utf-8', errors='replace').decode('utf-8')
        print(json.dumps({"success": False, "error": safe_err},
                         ensure_ascii=False))
        sys.exit(1)
