#!/usr/bin/env python3
"""
HTML-to-PPTX Converter
将 HTML 描述的幻灯片内容转换为 PPTX 演示文稿。
基于 Anthropic html2pptx 思想，使用 Python 实现原生转换。

支持的 HTML 结构：
  <section class="slide"> ... </section> — 每张幻灯片

支持的标签和样式：
  <h1>, <h2>, <h3> — 标题
  <p> — 段落
  <ul>, <li> — 列表
  <img> — 图片
  <table> — 表格
  class="cover" / "section" / "ending" — 特殊幻灯片类型
  data-kpi, data-value, data-trend — KPI 指示器
"""

import sys
import os
import re
import json
import tempfile
import shutil
from html.parser import HTMLParser
from xml.etree import ElementTree as ET

_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.dirname(_SCRIPT_DIR))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 中国商务配色
CLR = {
    "primary": "#1A3C6D", "accent": "#C0392B", "gold": "#8B6914",
    "text": "#2C2C2C", "text_sec": "#666666", "bg": "#F7F8FA",
    "bg_section": "#F0F3F7", "white": "#FFFFFF",
}


class SlideNode:
    def __init__(self):
        self.heading = ""
        self.subtitle = ""
        self.type = ""
        self.content = []
        self.presenter = ""


class ContentBlock:
    def __init__(self):
        self.tag = ""
        self.text = ""
        self.level = 0
        self.bold = False
        self.italic = False
        self.src = ""
        self.rows = []
        self.kpi_label = ""
        self.kpi_value = ""
        self.kpi_trend = ""


class HTMLToPPTXParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.slides = []
        self.current = SlideNode()
        self.content_block = None
        self.in_slide = False
        self.in_kpi = False
        self.in_table = False
        self.table_curr_row = []
        self.text_stack = []
        self.cover_title = ""
        self.cover_subtitle = ""
        self.presenter = ""
        self.style_tag = ""
        self.styles = {}

    def handle_starttag(self, tag, attrs):
        attrs_d = dict(attrs)
        cls = attrs_d.get("class", "")
        d_type = attrs_d.get("data-type", "")

        if tag == "style":
            self.style_tag = ""
        elif tag == "title" and not self.in_slide:
            pass
        elif tag == "section":
            self.in_slide = True
            self.current = SlideNode()
            self.current.type = cls if cls else d_type

        elif tag in ("h1", "h2", "h3"):
            self.content_block = ContentBlock()
            self.content_block.tag = tag
            self.content_block.bold = True

        elif tag == "p":
            self.content_block = ContentBlock()
            self.content_block.tag = "p"

        elif tag == "li":
            self.content_block = ContentBlock()
            self.content_block.tag = "li"

        elif tag == "img":
            self.content_block = ContentBlock()
            self.content_block.tag = "img"
            self.content_block.src = attrs_d.get("src", "")

        elif tag == "table":
            self.in_table = True
            self.content_block = ContentBlock()
            self.content_block.tag = "table"
            self.content_block.rows = []

        elif tag == "tr" and self.in_table:
            self.table_curr_row = []

        elif tag in ("td", "th") and self.in_table:
            pass

        elif tag == "div" and "kpi" in cls:
            self.in_kpi = True
            self.content_block = ContentBlock()
            self.content_block.tag = "kpi"
            self.content_block.kpi_label = attrs_d.get("data-label", "")
            self.content_block.kpi_value = attrs_d.get("data-value", "")
            self.content_block.kpi_trend = attrs_d.get("data-trend", "")

    def handle_endtag(self, tag):
        if tag == "section" and self.in_slide:
            self.in_slide = False
            self.slides.append(self.current)

        elif tag in ("h1", "h2", "h3", "p", "li", "img") and self.content_block:
            self.content_block.text = " ".join(self.text_stack).strip()
            self.text_stack = []
            if tag == "h1" and not self.in_slide:
                self.cover_title = self.content_block.text
            elif tag == "h2" and not self.in_slide:
                self.cover_subtitle = self.content_block.text
            elif self.in_slide:
                if tag == "h1":
                    self.current.heading = self.content_block.text
                elif tag == "h2" and self.current.type in ("cover", "cover-slide"):
                    self.current.subtitle = self.content_block.text
                else:
                    self.current.content.append(self.content_block)
            self.content_block = None

        elif tag == "div" and self.in_kpi:
            self.in_kpi = False
            self.content_block.text = " ".join(self.text_stack).strip()
            self.text_stack = []
            if self.in_slide:
                self.current.content.append(self.content_block)
            self.content_block = None

        elif tag in ("td", "th") and self.in_table:
            self.table_curr_row.append(" ".join(self.text_stack).strip())
            self.text_stack = []

        elif tag == "tr" and self.in_table:
            if self.content_block:
                self.content_block.rows.append(self.table_curr_row)

        elif tag == "table" and self.in_table:
            self.in_table = False
            if self.in_slide:
                self.current.content.append(self.content_block)
            self.content_block = None

        elif tag == "style":
            pass

    def handle_data(self, data):
        stripped = data.strip()
        if stripped:
            self.text_stack.append(stripped)


def _rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _font(run, size=16, bold=False, color=None):
    run.font.name = "Microsoft YaHei"
    rPr = run._element.get_or_add_rPr()
    ea = rPr.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", "Microsoft YaHei")
    rPr.append(ea)
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = _rgb(color)


def _textbox(slide, left, top, width, height, text, size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        _font(r, size, bold, color)
    return tb


def _accent_bar(slide, top=0, height=0.05, color="#C0392B"):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(top), Inches(13.333), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()


def convert_html_to_pptx(html_content, output_path):
    parser = HTMLToPPTXParser()
    parser.feed(html_content)
    parser.close()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = parser.slides
    if not slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _textbox(slide, 1, 3, 11, 1, "HTML 内容为空", size=24, bold=True, align=PP_ALIGN.CENTER)
        prs.save(output_path)
        return output_path

    total = len(slides)

    for idx, sn in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if sn.type in ("cover", "cover-slide"):
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = _rgb("#0D2137")
            _accent_bar(slide, 2.6, 0.05)
            _textbox(slide, 1.8, 1.8, 9.5, 2.4, sn.heading or parser.cover_title or "演示文稿",
                      size=44, bold=True, color="#FFFFFF")
            _textbox(slide, 1.8, 3.9, 9.5, 0.7, sn.subtitle or parser.cover_subtitle or "",
                      size=22, color="#E8C5C4")
            _textbox(slide, 1.8, 6.2, 9.5, 0.4, "内部资料 · 请注意保密",
                      size=10, color="#999999")

        elif sn.type == "section":
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = _rgb("#0D2137")
            _textbox(slide, 1.5, 2.2, 10, 1.2, ">>",
                      size=40, bold=True, color="#C0392B")
            _textbox(slide, 1.5, 3.6, 10, 0.8, "PART",
                      size=20, color="#999999")
            _textbox(slide, 1.5, 4.1, 10, 1.0, sn.heading,
                      size=32, bold=True, color="#FFFFFF")

        elif sn.type == "ending":
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = _rgb("#0D2137")
            _textbox(slide, 2.0, 2.5, 9.5, 1.2, "感谢聆听",
                      size=48, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)
            _accent_bar(slide, 3.6, 0.04)
            _textbox(slide, 2.0, 3.9, 9.5, 0.8, sn.heading,
                      size=20, color="#999999", align=PP_ALIGN.CENTER)
            _textbox(slide, 2.0, 5.0, 9.5, 0.6, "WebSQL AI · 智能数据分析平台",
                      size=14, color="#999999", align=PP_ALIGN.CENTER)

        else:
            # 标准内容页
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _rgb("#F7F8FA")
            _accent_bar(slide, 0, 0.05)

            if sn.heading:
                _textbox(slide, 1.2, 0.35, 10.5, 0.7, sn.heading,
                          size=26, bold=True, color="#1A3C6D")

            y = 1.5
            for block in sn.content:
                if block.tag == "p":
                    _textbox(slide, 1.2, y, 10.5, 0.5, block.text,
                              size=16, color="#2C2C2C")
                    y += 0.55

                elif block.tag == "h2":
                    _textbox(slide, 1.2, y, 10.5, 0.45, block.text,
                              size=22, bold=True, color="#1A3C6D")
                    y += 0.6

                elif block.tag == "h3":
                    _textbox(slide, 1.2, y, 10.5, 0.4, block.text,
                              size=18, bold=True, color="#666666")
                    y += 0.5

                elif block.tag == "li":
                    indent = block.level * 0.4
                    _textbox(slide, 1.5 + indent, y, 10 - indent, 0.35,
                              "\u2022 " + block.text,
                              size=15, color="#2C2C2C")
                    y += 0.4

                elif block.tag == "img" and block.src and os.path.exists(block.src):
                    slide.shapes.add_picture(block.src, Inches(1.2), Inches(y), Inches(10.5), Inches(4.0))
                    y += 4.2

                elif block.tag == "table" and block.rows:
                    rlen = len(block.rows)
                    clen = max(len(r) for r in block.rows[:rlen]) if block.rows else 1
                    if clen > 0:
                        max_r = min(rlen, 10)
                        tbl = slide.shapes.add_table(
                            max_r + 1, clen,
                            Inches(1.2), Inches(y),
                            Inches(10.5), Inches(0.35 * (max_r + 1))
                        ).table
                        for c in range(clen):
                            cell = tbl.cell(0, c)
                            cell.text = block.rows[0][c] if c < len(block.rows[0]) else ""
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _rgb("#1A3C6D")
                            for p in cell.text_frame.paragraphs:
                                for r in p.runs:
                                    _font(r, 10, True, "#FFFFFF")
                        for j in range(1, max_r + 1):
                            row_data = block.rows[j] if j < len(block.rows) else []
                            bg_color = "#FFFFFF" if j % 2 == 1 else "#F2F4F8"
                            for c in range(clen):
                                cell = tbl.cell(j, c)
                                cell.text = row_data[c] if c < len(row_data) else ""
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = _rgb(bg_color)
                                for p in cell.text_frame.paragraphs:
                                    for r in p.runs:
                                        _font(r, 9, False, "#2C2C2C")
                        y += 0.35 * (max_r + 1) + 0.3

                elif block.tag == "kpi":
                    kpi_text = ""
                    if block.kpi_label:
                        kpi_text += block.kpi_label
                    if block.kpi_value:
                        kpi_text += f": {block.kpi_value}"
                    if block.kpi_trend:
                        kpi_text += f"  {block.kpi_trend}"
                    _textbox(slide, 1.2, y, 10.5, 0.5, kpi_text,
                              size=20, bold=True, color="#C0392B")
                    y += 0.6

        _textbox(slide, 11.8, 7.08, 1.3, 0.3, f"{idx + 1} / {total}",
                  size=9, color="#999999", align=PP_ALIGN.RIGHT)

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        print("用法: python html2pptx.py <input.html> [output.pptx]", file=sys.stderr)
        print("       echo '<html>...' | python html2pptx.py --stdin -o output.pptx", file=sys.stderr)
        sys.exit(1)

    if sys.argv[1] == "--stdin":
        sys.stdin.reconfigure(encoding='utf-8')
        html_content = sys.stdin.read()
        output = "output.pptx"
        for i, arg in enumerate(sys.argv):
            if arg == "-o" and i + 1 < len(sys.argv):
                output = sys.argv[i + 1]
    else:
        with open(sys.argv[1], "r", encoding="utf-8") as f:
            html_content = f.read()
        output = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"

    sys.stdout.reconfigure(encoding='utf-8')
    result = convert_html_to_pptx(html_content, output)
    print(json.dumps({"success": True, "path": result}))
    return result


if __name__ == "__main__":
    main()
