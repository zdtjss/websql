"""
PPT 生成器 - 科技感数据分析风格（中国特色）
依赖: pip install python-pptx matplotlib numpy Pillow
用法: from tools.ppt_generator import PPTBuilder
"""

import os
import tempfile
from pathlib import Path

import matplotlib
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 中文字体配置
matplotlib.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'DejaVu Sans']
matplotlib.rcParams['axes.unicode_minus'] = False


# ═══════════════════════════════════════════════════════════════
# 科技感配色主题
# ═══════════════════════════════════════════════════════════════
class Theme:
    # PPT 颜色
    BG_DARK = RGBColor(0x0A, 0x16, 0x28)       # 深空蓝背景
    BG_CARD = RGBColor(0x12, 0x23, 0x3D)       # 卡片背景
    BLUE = RGBColor(0x00, 0xA8, 0xFF)          # 科技蓝
    CYAN = RGBColor(0x00, 0xF5, 0xD4)          # 电光青
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # 白色
    GRAY = RGBColor(0xB0, 0xC4, 0xD8)          # 浅灰蓝
    PURPLE = RGBColor(0x7B, 0x61, 0xFF)        # 数据紫
    ORANGE = RGBColor(0xFF, 0x6B, 0x35)        # 警示橙
    RED = RGBColor(0xE6, 0x39, 0x46)           # 中国红
    GOLD = RGBColor(0xFF, 0xD7, 0x00)          # 星光金
    GREEN = RGBColor(0x00, 0xE6, 0x96)         # 增长绿

    # matplotlib 图表配色
    CHART_COLORS = ['#00A8FF', '#00F5D4', '#7B61FF', '#FF6B35', '#E63946', '#FFD700', '#00E696']
    CHART_BG = '#0A1628'
    CHART_FACE = '#0F1E36'
    CHART_GRID = '#1E3A5F'


# ═══════════════════════════════════════════════════════════════
# PPTBuilder 主类
# ═══════════════════════════════════════════════════════════════
class PPTBuilder:
    def __init__(self, widescreen=True):
        self.prs = Presentation()
        if widescreen:
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        self._temp_files = []
        self._setup_matplotlib()

    def _setup_matplotlib(self):
        plt.style.use('dark_background')
        plt.rcParams.update({
            'figure.facecolor': Theme.CHART_FACE,
            'axes.facecolor': Theme.CHART_BG,
            'text.color': '#FFFFFF',
            'axes.labelcolor': '#B0C4D8',
            'xtick.color': '#B0C4D8',
            'ytick.color': '#B0C4D8',
            'axes.edgecolor': Theme.CHART_GRID,
            'grid.color': Theme.CHART_GRID,
            'grid.alpha': 0.3,
        })

    def _blank_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = Theme.BG_DARK
        return slide

    # ─── 基础元素 ───────────────────────────────────────────────

    def _add_text(self, slide, text, left, top, width, height,
                  size=Pt(16), color=Theme.GRAY, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = size
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = align
            p.space_after = Pt(6)
        return box

    def _add_line(self, slide, left, top, width, color=Theme.BLUE):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

    def _add_card(self, slide, left, top, width, height, border_color=Theme.BLUE):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        s.fill.solid()
        s.fill.fore_color.rgb = Theme.BG_CARD
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
        return s

    def _save_chart(self, fig):
        path = os.path.join(tempfile.gettempdir(), f'ppt_chart_{len(self._temp_files)}.png')
        fig.savefig(path, dpi=150, bbox_inches='tight', facecolor=Theme.CHART_FACE, edgecolor='none')
        plt.close(fig)
        self._temp_files.append(path)
        return path

    # ─── 页面类型 ───────────────────────────────────────────────

    def add_cover(self, title, subtitle="", date="", author=""):
        """封面页"""
        slide = self._blank_slide()
        # 左侧发光线
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Pt(4), Inches(7.5))
        s.fill.solid(); s.fill.fore_color.rgb = Theme.BLUE; s.line.fill.background()
        # 底部线
        self._add_line(slide, Inches(1.2), Inches(6.8), Inches(11), Theme.CYAN)
        # 文字
        self._add_text(slide, title, Inches(1.5), Inches(2), Inches(10), Inches(1.5),
                       size=Pt(44), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(1.5), Inches(3.7), Inches(3))
        if subtitle:
            self._add_text(slide, subtitle, Inches(1.5), Inches(4.0), Inches(8), Inches(0.8),
                           size=Pt(20), color=Theme.CYAN)
        bottom = f"{date}    {author}".strip()
        if bottom:
            self._add_text(slide, bottom, Inches(1.5), Inches(5.5), Inches(6), Inches(0.6),
                           size=Pt(14), color=Theme.GRAY)

    def add_toc(self, sections):
        """目录页"""
        slide = self._blank_slide()
        self._add_text(slide, "目录", Inches(0.8), Inches(0.4), Inches(4), Inches(0.8),
                       size=Pt(32), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(0.8), Inches(1.1), Inches(2))
        for i, sec in enumerate(sections):
            y = Inches(1.8) + Inches(i * 0.7)
            num = f"0{i+1}" if i < 9 else str(i+1)
            self._add_text(slide, num, Inches(1.0), y, Inches(0.8), Inches(0.6),
                           size=Pt(24), color=Theme.BLUE, bold=True)
            self._add_text(slide, sec, Inches(1.9), y, Inches(8), Inches(0.6),
                           size=Pt(18), color=Theme.WHITE)

    def add_section_divider(self, title, desc=""):
        """章节分隔页"""
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(2), Inches(2.5), Inches(9), Inches(1.2),
                       size=Pt(36), color=Theme.WHITE, bold=True, align=PP_ALIGN.CENTER)
        self._add_line(slide, Inches(5.5), Inches(3.8), Inches(2.3), Theme.CYAN)
        if desc:
            self._add_text(slide, desc, Inches(2), Inches(4.2), Inches(9), Inches(1),
                           size=Pt(16), color=Theme.GRAY, align=PP_ALIGN.CENTER)

    def add_kpi_page(self, title, kpis):
        """KPI 大数字展示页，kpis: list of {label, value, change, trend}"""
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(0.8), Inches(0.4), Inches(6), Inches(0.7),
                       size=Pt(28), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(0.8), Inches(1.0), Inches(2))

        n = len(kpis)
        card_w = Inches(min(2.8, 11.0 / n))
        gap = Inches(0.3)
        start_x = Inches(0.8)

        for i, kpi in enumerate(kpis):
            x = start_x + i * (card_w + gap)
            y = Inches(1.8)
            self._add_card(slide, x, y, card_w, Inches(4.5))
            # 指标名
            self._add_text(slide, kpi['label'], x + Inches(0.3), y + Inches(0.4),
                           card_w - Inches(0.6), Inches(0.5), size=Pt(13), color=Theme.GRAY)
            # 大数字
            self._add_text(slide, kpi['value'], x + Inches(0.3), y + Inches(1.2),
                           card_w - Inches(0.6), Inches(1.2), size=Pt(28), color=Theme.WHITE, bold=True)
            # 变化
            trend_color = Theme.GREEN if kpi.get('trend') == 'up' else Theme.RED if kpi.get('trend') == 'down' else Theme.GRAY
            arrow = "▲" if kpi.get('trend') == 'up' else "▼" if kpi.get('trend') == 'down' else "─"
            self._add_text(slide, f"{arrow} {kpi.get('change', '')}", x + Inches(0.3), y + Inches(2.8),
                           card_w - Inches(0.6), Inches(0.6), size=Pt(14), color=trend_color)

    def add_chart_page(self, title, chart_type, data, insights=None, layout="left_chart"):
        """图表+文字页"""
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                       size=Pt(26), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(0.8), Inches(0.95), Inches(2))

        fig = self._create_chart(chart_type, data)
        img_path = self._save_chart(fig)

        if layout == "left_chart":
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4), Inches(7.5), Inches(5.5))
            if insights:
                text = '\n'.join(f"▸ {s}" for s in insights)
                self._add_card(slide, Inches(8.3), Inches(1.6), Inches(4.5), Inches(5))
                self._add_text(slide, text, Inches(8.6), Inches(2.0), Inches(4), Inches(4.5),
                               size=Pt(14), color=Theme.GRAY)
        elif layout == "full_chart":
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
            if insights:
                text = '  |  '.join(insights)
                self._add_text(slide, text, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
                               size=Pt(12), color=Theme.CYAN)
        elif layout == "top_chart":
            slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.3), Inches(11.5), Inches(4))
            if insights:
                text = '\n'.join(f"▸ {s}" for s in insights)
                self._add_text(slide, text, Inches(1.0), Inches(5.5), Inches(11), Inches(1.8),
                               size=Pt(14), color=Theme.GRAY)

    def add_dual_chart_page(self, title, chart1_config, chart2_config):
        """双图表页"""
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                       size=Pt(26), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(0.8), Inches(0.95), Inches(2))

        fig1 = self._create_chart(chart1_config['type'], chart1_config['data'])
        fig2 = self._create_chart(chart2_config['type'], chart2_config['data'])
        p1 = self._save_chart(fig1)
        p2 = self._save_chart(fig2)

        slide.shapes.add_picture(p1, Inches(0.3), Inches(1.4), Inches(6.2), Inches(5.5))
        slide.shapes.add_picture(p2, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.5))

    def add_text_page(self, title, content, highlight_indices=None):
        """纯文字要点页"""
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                       size=Pt(28), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(0.8), Inches(1.05), Inches(2))

        highlight_indices = highlight_indices or []
        for i, item in enumerate(content):
            y = Inches(1.6) + Inches(i * 0.75)
            color = Theme.CYAN if i in highlight_indices else Theme.GRAY
            self._add_text(slide, f"▸ {item}", Inches(1.2), y, Inches(11), Inches(0.7),
                           size=Pt(16), color=color)

    def add_comparison_page(self, title, left_title, left_items, right_title, right_items):
        """左右对比页"""
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                       size=Pt(28), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(0.8), Inches(1.05), Inches(2))

        # 左卡片
        self._add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), Theme.BLUE)
        self._add_text(slide, left_title, Inches(0.8), Inches(1.7), Inches(5), Inches(0.6),
                       size=Pt(18), color=Theme.BLUE, bold=True)
        left_text = '\n'.join(f"• {x}" for x in left_items)
        self._add_text(slide, left_text, Inches(1.0), Inches(2.5), Inches(5), Inches(4),
                       size=Pt(14), color=Theme.GRAY)

        # 右卡片
        self._add_card(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), Theme.CYAN)
        self._add_text(slide, right_title, Inches(7.1), Inches(1.7), Inches(5), Inches(0.6),
                       size=Pt(18), color=Theme.CYAN, bold=True)
        right_text = '\n'.join(f"• {x}" for x in right_items)
        self._add_text(slide, right_text, Inches(7.3), Inches(2.5), Inches(5), Inches(4),
                       size=Pt(14), color=Theme.GRAY)

    def add_summary_page(self, title, conclusions):
        """总结页"""
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                       size=Pt(30), color=Theme.WHITE, bold=True)
        self._add_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), Theme.GOLD)

        for i, c in enumerate(conclusions):
            y = Inches(1.8) + Inches(i * 0.9)
            num_color = Theme.GOLD
            self._add_text(slide, str(i+1), Inches(1.0), y, Inches(0.5), Inches(0.6),
                           size=Pt(22), color=num_color, bold=True)
            self._add_text(slide, c, Inches(1.7), y, Inches(10), Inches(0.7),
                           size=Pt(16), color=Theme.WHITE)

    def add_thank_you(self, text="感谢聆听", contact=""):
        """致谢页"""
        slide = self._blank_slide()
        self._add_text(slide, text, Inches(2), Inches(2.8), Inches(9), Inches(1.2),
                       size=Pt(40), color=Theme.WHITE, bold=True, align=PP_ALIGN.CENTER)
        self._add_line(slide, Inches(5.5), Inches(4.2), Inches(2.3), Theme.CYAN)
        if contact:
            self._add_text(slide, contact, Inches(2), Inches(4.8), Inches(9), Inches(0.8),
                           size=Pt(14), color=Theme.GRAY, align=PP_ALIGN.CENTER)

    # ─── 图表生成 ───────────────────────────────────────────────

    def _create_chart(self, chart_type, data):
        """根据类型生成 matplotlib 图表"""
        creators = {
            'line': self._chart_line,
            'bar': self._chart_bar,
            'horizontal_bar': self._chart_hbar,
            'pie': self._chart_pie,
            'donut': self._chart_donut,
            'scatter': self._chart_scatter,
            'radar': self._chart_radar,
            'heatmap': self._chart_heatmap,
            'area': self._chart_area,
            'stacked_bar': self._chart_stacked_bar,
        }
        fn = creators.get(chart_type, self._chart_bar)
        return fn(data)

    def _chart_line(self, data):
        fig, ax = plt.subplots(figsize=(8, 5))
        for i, s in enumerate(data['series']):
            ax.plot(data['categories'], s['values'], marker='o', linewidth=2,
                    color=Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)], label=s['name'])
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        ax.legend(loc='upper left', framealpha=0.3)
        ax.grid(True, alpha=0.2)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        fig.tight_layout()
        return fig

    def _chart_bar(self, data):
        fig, ax = plt.subplots(figsize=(8, 5))
        cats = data['categories']
        n_series = len(data['series'])
        width = 0.7 / n_series
        x = np.arange(len(cats))
        for i, s in enumerate(data['series']):
            offset = (i - n_series/2 + 0.5) * width
            bars = ax.bar(x + offset, s['values'], width, label=s['name'],
                          color=Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)], alpha=0.9)
            for bar, val in zip(bars, s['values']):
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                        str(val), ha='center', va='bottom', fontsize=9, color='#00F5D4')
        ax.set_xticks(x)
        ax.set_xticklabels(cats)
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        if n_series > 1:
            ax.legend(framealpha=0.3)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(axis='y', alpha=0.2)
        fig.tight_layout()
        return fig

    def _chart_hbar(self, data):
        fig, ax = plt.subplots(figsize=(8, 5))
        cats = data['categories']
        values = data['series'][0]['values']
        colors = [Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)] for i in range(len(cats))]
        ax.barh(cats, values, color=colors, alpha=0.9)
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        fig.tight_layout()
        return fig

    def _chart_pie(self, data):
        fig, ax = plt.subplots(figsize=(6, 6))
        colors = Theme.CHART_COLORS[:len(data['labels'])]
        ax.pie(data['values'], labels=data['labels'], colors=colors, autopct='%1.1f%%',
               textprops={'color': 'white', 'fontsize': 11}, startangle=90)
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        fig.tight_layout()
        return fig

    def _chart_donut(self, data):
        fig, ax = plt.subplots(figsize=(6, 6))
        colors = Theme.CHART_COLORS[:len(data['labels'])]
        wedges, texts, autotexts = ax.pie(
            data['values'], labels=data['labels'], colors=colors, autopct='%1.1f%%',
            textprops={'color': 'white', 'fontsize': 11}, startangle=90,
            pctdistance=0.8, wedgeprops={'width': 0.4})
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        fig.tight_layout()
        return fig

    def _chart_scatter(self, data):
        fig, ax = plt.subplots(figsize=(8, 5))
        ax.scatter(data['x'], data['y'], c=Theme.CHART_COLORS[0], alpha=0.7, s=60, edgecolors='white', linewidth=0.5)
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        ax.set_xlabel(data.get('x_label', ''))
        ax.set_ylabel(data.get('y_label', ''))
        ax.grid(True, alpha=0.2)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        fig.tight_layout()
        return fig

    def _chart_radar(self, data):
        fig, ax = plt.subplots(figsize=(6, 6), subplot_kw=dict(polar=True))
        cats = data['categories']
        n = len(cats)
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        angles += angles[:1]
        for i, s in enumerate(data['series']):
            vals = s['values'] + s['values'][:1]
            ax.plot(angles, vals, 'o-', linewidth=2,
                    color=Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)], label=s['name'])
            ax.fill(angles, vals, alpha=0.15, color=Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)])
        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(cats, fontsize=10)
        ax.set_title(data.get('title', ''), fontsize=14, pad=20, color='white')
        ax.legend(loc='upper right', framealpha=0.3)
        fig.tight_layout()
        return fig

    def _chart_heatmap(self, data):
        fig, ax = plt.subplots(figsize=(8, 5))
        values = np.array(data['values'])
        im = ax.imshow(values, cmap='YlOrRd', aspect='auto')
        ax.set_xticks(range(len(data['x_labels'])))
        ax.set_xticklabels(data['x_labels'], fontsize=9)
        ax.set_yticks(range(len(data['y_labels'])))
        ax.set_yticklabels(data['y_labels'], fontsize=9)
        fig.colorbar(im, ax=ax, shrink=0.8)
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        fig.tight_layout()
        return fig

    def _chart_area(self, data):
        fig, ax = plt.subplots(figsize=(8, 5))
        for i, s in enumerate(data['series']):
            ax.fill_between(data['categories'], s['values'], alpha=0.3,
                            color=Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)])
            ax.plot(data['categories'], s['values'], linewidth=2,
                    color=Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)], label=s['name'])
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        ax.legend(framealpha=0.3)
        ax.grid(True, alpha=0.2)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        fig.tight_layout()
        return fig

    def _chart_stacked_bar(self, data):
        fig, ax = plt.subplots(figsize=(8, 5))
        cats = data['categories']
        x = np.arange(len(cats))
        bottom = np.zeros(len(cats))
        for i, s in enumerate(data['series']):
            ax.bar(x, s['values'], bottom=bottom, label=s['name'],
                   color=Theme.CHART_COLORS[i % len(Theme.CHART_COLORS)], alpha=0.9)
            bottom += np.array(s['values'])
        ax.set_xticks(x)
        ax.set_xticklabels(cats)
        ax.set_title(data.get('title', ''), fontsize=14, pad=10, color='white')
        ax.legend(framealpha=0.3)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        fig.tight_layout()
        return fig

    # ─── 保存 ──────────────────────────────────────────────────

    def save(self, filepath):
        """保存 PPT 并清理临时文件"""
        Path(filepath).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(filepath)
        for f in self._temp_files:
            try:
                os.remove(f)
            except OSError:
                pass
        self._temp_files.clear()
        print(f"✅ PPT 已生成: {filepath}")
