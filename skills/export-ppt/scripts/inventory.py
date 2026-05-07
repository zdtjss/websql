#!/usr/bin/env python3
"""
PPT Template Inventory — 提取模板中所有形状及其属性
输出 JSON 格式的仓库清单，包含占位符类型、位置、字号、文本等信息。
基于 Anthropic pptx skill 的 inventory.py 思想，适配 Python 环境。
"""

import sys
import os
import json
from pptx import Presentation
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.util import Inches, Emu


_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.dirname(_SCRIPT_DIR))


def extract_inventory(pptx_path, output_path=None):
    prs = Presentation(pptx_path)
    inventory = {}

    for slide_idx, slide in enumerate(prs.slides):
        slide_key = f"slide-{slide_idx}"
        shapes_data = []
        shape_idx = 0

        shapes_sorted = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        for shape in shapes_sorted:
            if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                continue

            pt = None
            if shape.is_placeholder:
                ph = shape.placeholder_format
                if ph.type is not None:
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    ph_map = {
                        PP_PLACEHOLDER.TITLE: "TITLE",
                        PP_PLACEHOLDER.CENTER_TITLE: "CENTER_TITLE",
                        PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
                        PP_PLACEHOLDER.BODY: "BODY",
                        PP_PLACEHOLDER.OBJECT: "OBJECT",
                        PP_PLACEHOLDER.SLIDE_NUMBER: "SLIDE_NUMBER",
                    }
                    pt = ph_map.get(ph.type, str(ph.type))

            if pt == "SLIDE_NUMBER":
                continue

            shape_entry = {
                "left": round(shape.left / 914400, 2) if shape.left else 0,
                "top": round(shape.top / 914400, 2) if shape.top else 0,
                "width": round(shape.width / 914400, 2) if shape.width else 0,
                "height": round(shape.height / 914400, 2) if shape.height else 0,
                "placeholder_type": pt,
                "paragraphs": [],
            }

            for para in shape.text_frame.paragraphs:
                if not para.text.strip():
                    continue
                para_data = {
                    "text": para.text,
                }

                if para.alignment is not None:
                    from pptx.enum.text import PP_ALIGN
                    am = {PP_ALIGN.LEFT: "LEFT", PP_ALIGN.CENTER: "CENTER", PP_ALIGN.RIGHT: "RIGHT"}
                    if para.alignment in am:
                        para_data["alignment"] = am[para.alignment]

                for run in para.runs:
                    if run.font.size:
                        para_data["font_size"] = run.font.size / 12700
                    if run.font.bold:
                        para_data["bold"] = run.font.bold
                    if run.font.italic:
                        para_data["italic"] = run.font.italic
                    if run.font.name:
                        para_data["font_name"] = run.font.name
                    if run.font.color and run.font.color.rgb:
                        para_data["color"] = str(run.font.color.rgb)
                    break

                shape_entry["paragraphs"].append(para_data)

            if shape_entry["paragraphs"]:
                shapes_data.append(shape_entry)
                shape_idx += 1

        inventory[slide_key] = {}
        for i, sd in enumerate(shapes_data):
            inventory[slide_key][f"shape-{i}"] = sd

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(inventory, f, ensure_ascii=False, indent=2)
        print(f"已提取 {len(inventory)} 张幻灯片的文字 -> {output_path}")
    else:
        print(json.dumps(inventory, ensure_ascii=False, indent=2))

    return inventory


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python inventory.py <template.pptx> [output.json]", file=sys.stderr)
        sys.exit(1)
    pptx_path = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else None
    extract_inventory(pptx_path, output)
