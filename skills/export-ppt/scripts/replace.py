#!/usr/bin/env python3
"""
PPT Template Replace — 根据替换 JSON 批量更新幻灯片文字
保留原有格式（字体、颜色、对齐），只替换文字内容和样式覆盖。
基于 Anthropic pptx skill 的 replace.py 思想。
"""

import sys
import os
import json
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.dirname(_SCRIPT_DIR))


def replace_text(pptx_path, replacement_json_path, output_path):
    prs = Presentation(pptx_path)

    with open(replacement_json_path, "r", encoding="utf-8") as f:
        replacements = json.load(f)

    for slide_idx, slide in enumerate(prs.slides):
        slide_key = f"slide-{slide_idx}"
        if slide_key not in replacements:
            continue

        shapes_sorted = sorted(
            [s for s in slide.shapes if hasattr(s, "text_frame") and s.has_text_frame],
            key=lambda s: (s.top or 0, s.left or 0)
        )

        shape_list = []
        for shape in shapes_sorted:
            if shape.is_placeholder:
                from pptx.enum.shapes import PP_PLACEHOLDER
                if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                    continue
            shape_list.append(shape)

        slide_replacements = replacements[slide_key]

        for shape_key, shape_data in slide_replacements.items():
            shape_num = int(shape_key.split("-")[1])
            if shape_num >= len(shape_list):
                continue

            shape = shape_list[shape_num]
            tf = shape.text_frame
            tf.clear()

            paragraphs = shape_data.get("paragraphs", [])
            for i, para_data in enumerate(paragraphs):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                if "alignment" in para_data:
                    am = {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT}
                    p.alignment = am.get(para_data["alignment"], PP_ALIGN.LEFT)

                text = para_data.get("text", "")
                run = p.add_run()
                run.text = text

                if "font_size" in para_data:
                    run.font.size = Pt(para_data["font_size"])
                if "bold" in para_data:
                    run.font.bold = para_data["bold"]
                if "italic" in para_data:
                    run.font.italic = para_data["italic"]
                if "font_name" in para_data:
                    run.font.name = para_data["font_name"]
                if "color" in para_data:
                    c = para_data["color"]
                    if len(c) == 6:
                        run.font.color.rgb = RGBColor(
                            int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))

                if "bullet" in para_data and para_data["bullet"]:
                    p.level = para_data.get("level", 0)
                    from pptx.oxml.ns import qn
                    buChar = p._pPr.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
                    buChar.set("char", "\u2022")
                    existing_bu = p._pPr.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
                    for eb in existing_bu:
                        p._pPr.remove(eb)
                    p._pPr.append(buChar)

    prs.save(output_path)
    print(f"文字替换完成 -> {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("用法: python replace.py <template.pptx> <replacement.json> <output.pptx>", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    replacement_json = sys.argv[2]
    output_path = sys.argv[3]
    replace_text(pptx_path, replacement_json, output_path)
