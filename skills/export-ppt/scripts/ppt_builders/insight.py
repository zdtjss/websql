"""
核心发现页构建器
生成要点分析、结论建议等内容的幻灯片，支持 ▶ 前缀要点展示。
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class InsightBuilder:
    """核心发现与建议构建器"""

    def __init__(self, template):
        self.template = template
        self.scheme = template.scheme

    def build(self, prs, highlights, slide_num, total):
        from shared.utils import clean_surrogates

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.template._solid_bg(slide, self.scheme["bg_slide"])
        self.template._accent_bar(slide, 0, 0.05, self.scheme["accent"])
        self.template._textbox(slide, 1.2, 0.35, 10.0, 0.8, "核心发现与建议",
                               size=30, bold=True, color=self.scheme["primary"])

        if not highlights:
            highlights = ["本次分析未发现显著异常模式。"]

        max_items = self.template.config.get("limits", "ppt_max_highlights")
        text = "\n\n".join(f"▶  {clean_surrogates(h)}" for h in highlights[:max_items])
        self.template._textbox(slide, 1.2, 1.6, 10.5, 5.2, text,
                               size=18, color=self.scheme["text_primary"])
        self.template._page_num(slide, slide_num, total)
        return slide

    def build_smart(self, prs, findings_with_actions, slide_num, total):
        """带行动建议的结构化发现"""
        from shared.utils import clean_surrogates
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.template._solid_bg(slide, self.scheme["bg_slide"])
        self.template._accent_bar(slide, 0, 0.05, self.scheme["accent"])
        self.template._textbox(slide, 1.2, 0.35, 10.0, 0.8, "核心发现与行动建议",
                               size=30, bold=True, color=self.scheme["primary"])

        y = 1.5
        for i, item in enumerate(findings_with_actions[:6]):
            finding = item.get("finding", "")
            action = item.get("action", "")
            metric = item.get("metric", "")

            header = f"发现 {i + 1}：" + (f" [指标：{metric}]" if metric else "")
            self.template._textbox(slide, 1.2, y, 10.5, 0.35,
                                   header, size=16, bold=True, color=self.scheme["accent"])
            y += 0.38
            self.template._textbox(slide, 1.5, y, 10.0, 0.3,
                                   clean_surrogates(finding), size=13, color=self.scheme["text_primary"])
            y += 0.3
            if action:
                self.template._textbox(slide, 1.5, y, 10.0, 0.25,
                                       f"▹ 建议：{clean_surrogates(action)}", size=12,
                                       color=self.scheme["text_secondary"])
                y += 0.28
            y += 0.15

        self.template._page_num(slide, slide_num, total)
        return slide
