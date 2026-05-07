"""
数据页构建器
生成数据全景/汇总、数据明细表格等幻灯片，支持统计分析和亮点提炼。
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class DataBuilder:
    """数据内容构建器"""

    def __init__(self, template):
        self.template = template
        self.scheme = template.scheme
        self.config = template.config

    def build_summary(self, prs, qr_summary, numeric_cols, slide_num, total):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.template._solid_bg(slide, self.scheme["bg_slide"])
        self.template._accent_bar(slide, 0, 0.05, self.scheme["accent"])
        self.template._textbox(slide, 1.2, 0.35, 10.0, 0.8, "数据全景",
                               size=30, bold=True, color=self.scheme["primary"])
        from shared.utils import clean_surrogates, format_number_cn

        total_rows = qr_summary.get("totalRows", 0)
        total_cols = qr_summary.get("totalCols", 0)
        lines = [f"▸  数据规模：{format_number_cn(total_rows)} 条记录  ·  {total_cols} 个维度"]
        if numeric_cols:
            cols_text = clean_surrogates("、".join(numeric_cols[:6]))
            lines.append(f"▸  核心指标：{cols_text}")

        stats = qr_summary.get("stats", {})
        stat_lines = []
        for col, info in list(stats.items())[:8]:
            col_clean = clean_surrogates(col)
            stat_lines.append(f"◦  {col_clean}  |  均值 {info['avg']:,.2f}  |  区间 [{info['min']:,.2f}, {info['max']:,.2f}]")

        if stat_lines:
            lines.append("")
            lines.append("▹  指标统计：")
            lines.extend(stat_lines)

        self.template._textbox(slide, 1.2, 1.6, 10.5, 5.0, "\n".join(lines),
                               size=16, color=self.scheme["text_primary"])
        self.template._page_num(slide, slide_num, total)
        return slide

    def build_table(self, prs, title, columns, data_rows, slide_num, total):
        from shared.utils import clean_surrogates, truncate_text
        max_cols = self.config.get("limits", "ppt_max_table_columns")
        max_rows = self.config.get("limits", "ppt_max_table_rows")

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.template._solid_bg(slide, self.scheme["bg_slide"])
        self.template._accent_bar(slide, 0, 0.05, self.scheme["accent"])
        self.template._textbox(slide, 1.2, 0.3, 10.0, 0.7, clean_surrogates(title),
                               size=26, bold=True, color=self.scheme["primary"])

        dc = min(len(columns), max_cols)
        dr = min(len(data_rows), max_rows)
        rows_total = dr + 1
        row_h = 0.35

        tbl = slide.shapes.add_table(rows_total, dc,
                                     Inches(1.2), Inches(1.2),
                                     Inches(10.5), Inches(row_h * rows_total)).table

        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor

        for c in range(dc):
            cell = tbl.cell(0, c)
            col_name = truncate_text(columns[c], 16)
            cell.text = col_name
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    self.template._font(run, size=10, bold=True, color=self.scheme["bg_white"])
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.template._rgb(self.scheme["table_header"])

        for r in range(dr):
            bg = self.scheme["bg_white"] if r % 2 == 0 else self.scheme["table_stripe"]
            for c in range(dc):
                cell = tbl.cell(r + 1, c)
                val = truncate_text(str(data_rows[r].get(columns[c], "")), 24)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        self.template._font(run, size=9.5)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.template._rgb(bg)

        if len(data_rows) > dr:
            self.template._textbox(slide, 1.2, 1.2 + row_h * rows_total + 0.15, 10.5, 0.3,
                                   f"※ 展示 {dr}/{len(data_rows)} 条", size=11, color=self.scheme["text_muted"])

        self.template._page_num(slide, slide_num, total)
        return slide
