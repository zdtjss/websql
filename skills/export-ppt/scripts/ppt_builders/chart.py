"""
图表页构建器
在幻灯片中嵌入 PNG 图表图片，支持多图展示。
"""

import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


class ChartBuilder:
    """图表幻灯片构建器"""

    def __init__(self, template):
        self.template = template
        self.scheme = template.scheme

    def build_single(self, prs, chart_title, chart_path, slide_num, total):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.template._solid_bg(slide, self.scheme["bg_slide"])
        self.template._accent_bar(slide, 0, 0.05, self.scheme["accent"])
        self.template._textbox(slide, 1.2, 0.35, 10.5, 0.7, chart_title,
                               size=26, bold=True, color=self.scheme["primary"])

        if os.path.exists(chart_path):
            slide.shapes.add_picture(chart_path, Inches(1.2), Inches(1.4),
                                     Inches(10.5), Inches(5.2))

        self.template._page_num(slide, slide_num, total)
        return slide

    def build_multi(self, prs, chart_paths, slide_num, total):
        slides = []
        for i, cp in enumerate(chart_paths):
            if os.path.exists(cp):
                fn = os.path.basename(cp)
                title = fn.replace(".png", "").replace("_", " ")
                s = self.build_single(prs, title, cp, slide_num + i, total)
                slides.append(s)
        return slides

    def build_comparison(self, prs, title, charts_list, slide_num, total):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.template._solid_bg(slide, self.scheme["bg_slide"])
        self.template._accent_bar(slide, 0, 0.05, self.scheme["accent"])
        self.template._textbox(slide, 1.2, 0.3, 10.0, 0.7, title,
                               size=24, bold=True, color=self.scheme["primary"])

        n = len(charts_list)
        w_per = min(10.0 / n, 5.0)
        for i, (path, label) in enumerate(charts_list):
            x = 1.2 + i * w_per + (w_per - 4.5) / 2
            if os.path.exists(path):
                slide.shapes.add_picture(path, Inches(x), Inches(1.4),
                                         Inches(min(4.5, w_per)), Inches(3.2))
                self.template._textbox(slide, x, 4.7, 4.5, 0.3, label,
                                       size=11, color=self.scheme["text_secondary"], align=PP_ALIGN.CENTER)

        self.template._page_num(slide, slide_num, total)
        return slide
