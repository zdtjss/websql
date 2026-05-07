"""
PPT 基础模板
定义所有幻灯片模板的通用接口和共享渲染逻辑。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


class PPTTemplate:
    name = "base"
    label = "基础模板"

    def __init__(self, scheme=None):
        from shared.colors import ColorPalette
        from shared.fonts import FontManager
        from shared.config import SkillConfig

        self.config = SkillConfig()
        self.scheme = scheme or ColorPalette.get("chinese_business")
        self.fonts = FontManager

        self.slide_width = Inches(self.config.get("layout", "ppt", "slide_width_inches"))
        self.slide_height = Inches(self.config.get("layout", "ppt", "slide_height_inches"))
        self.margin = self.config.get("layout", "ppt", "margin_left")
        self.accent_height = self.config.get("layout", "ppt", "accent_bar_height")

        self._rgb_cache = {}

    def _rgb(self, hex_color):
        if hex_color not in self._rgb_cache:
            h = hex_color.lstrip("#")
            self._rgb_cache[hex_color] = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        return self._rgb_cache[hex_color]

    def _font(self, run, size=18, bold=False, color=None, font_cn=None):
        from shared.utils import clean_surrogates
        fn_cn = font_cn or self.fonts.get_cn_font()
        run.font.name = fn_cn
        el = run.font._element
        el.get_or_add_latin().set("typeface", fn_cn)
        ea = el.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
        ea.set("typeface", fn_cn)
        el.append(ea)
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = self._rgb(color) if isinstance(color, str) else color

    def _gradient_bg(self, slide, c1, c2, c3):
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_angle = 90.0
        gsLst = fill._fill._gradFill.get_or_add_gsLst()
        gs_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        while len(gsLst) < 3:
            new_gs = etree.SubElement(gsLst, f"{gs_ns}gs")
            new_gs.set("pos", "100000")
            srgb = etree.SubElement(new_gs, f"{gs_ns}srgbClr")
            srgb.set("val", "000000")

        def to_rgb(c):
            return self._rgb(c) if isinstance(c, str) else c

        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[0].color.rgb = to_rgb(c1)
        fill.gradient_stops[1].position = 0.5
        fill.gradient_stops[1].color.rgb = to_rgb(c2)
        fill.gradient_stops[2].position = 1.0
        fill.gradient_stops[2].color.rgb = to_rgb(c3)

    def _solid_bg(self, slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(color) if isinstance(color, str) else color

    def _accent_bar(self, slide, top=0, height=None, color=None, width=None):
        h = height or self.accent_height
        c = color or self.scheme["accent"]
        w = width or self.slide_width
        shape = slide.shapes.add_shape(1, Inches(0), Inches(top), w, Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(c) if isinstance(c, str) else c
        shape.line.fill.background()

    def _left_bar(self, slide, left=0.8, top=1.2, height=5.6, width=0.06, color=None):
        c = color or self.scheme["accent"]
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(c) if isinstance(c, str) else c
        shape.line.fill.background()

    def _textbox(self, slide, left, top, width, height, text, size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, font_cn=None):
        from shared.utils import clean_surrogates
        text = clean_surrogates(text)
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        tb.text_frame.word_wrap = True
        tf = tb.text_frame
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            self._font(r, size, bold, color, font_cn)
        return tb

    def _page_num(self, slide, n, total):
        self._textbox(slide, 11.8, 7.08, 1.3, 0.3, f"{n} / {total}",
                      size=9, color=self.scheme["text_muted"], align=PP_ALIGN.RIGHT)

    def create_presentation(self):
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        return prs
