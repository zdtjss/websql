# -*- coding: utf-8 -*-
"""build_templates.py — 生成 export-ppt / export-word 的基线模板文件（一次性脚本）。

用法：
    python skills/lib/build_templates.py

产物：
  - skills/export-ppt/templates/slides_template.pptx  （16:9 深色主题，母版 + 版式）
  - skills/export-word/templates/report_template.docx （封面 / 目录 / 样式集 / 页眉页脚）

换肤 / 换版式方式（模板驱动的核心收益）：
  1. 直接用 PowerPoint / WPS / Word 编辑模板文件（推荐，零代码）
  2. 或修改本脚本后重新生成

版式索引约定（export_ppt.py 依赖，改动需同步脚本）：
  layout[0] 封面   layout[1] 内容/目录/表格/KPI   layout[2] 章节分隔
  layout[5] 图表   layout[6] 致谢（空白版式）
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn as dqn
from docx.shared import Pt as DPt
from docx.shared import RGBColor as DRGBColor

# ═══════════════════════════════════════════════════════════════
# 品牌色（与 export_ppt.py 的 Theme 保持一致）
# ═══════════════════════════════════════════════════════════════
BG_DARK = RGBColor(0x0A, 0x16, 0x28)      # 深空蓝背景
BG_CARD = RGBColor(0x12, 0x23, 0x3D)      # 卡片背景
BLUE = RGBColor(0x00, 0xA8, 0xFF)         # 科技蓝
CYAN = RGBColor(0x00, 0xF5, 0xD4)         # 电光青
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xB0, 0xC4, 0xD8)         # 浅灰蓝
GOLD = RGBColor(0xFF, 0xD7, 0x00)         # 星光金

EA_FONT = "微软雅黑"


def _set_ea_font(rPr, typeface):
    """在 a:latin 之后追加/更新 a:ea 东亚字体。"""
    latin = rPr.find(qn("a:latin"))
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    ea.set("typeface", typeface)


def _style_placeholder(ph, size, color, bold, align=None):
    """设置版式占位符的默认文字样式（段落级 defRPr，克隆到幻灯片后继承）。"""
    for para in ph.text_frame.paragraphs:
        para.font.size = Pt(size)
        para.font.color.rgb = color
        para.font.bold = bold
        if align is not None:
            para.alignment = align
        pPr = para._p.get_or_add_pPr()
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = pPr.makeelement(qn("a:defRPr"), {})
            pPr.append(defRPr)
        _set_ea_font(defRPr, EA_FONT)


def _move_ph(ph, left, top, width, height):
    ph.left, ph.top, ph.width, ph.height = (
        Inches(left), Inches(top), Inches(width), Inches(height))


def _add_bar(owner, left, top, width, height, color, name=None):
    """装饰矩形（XML 方式添加，母版/版式/幻灯片通用）。

    python-pptx 的 MasterShapes / LayoutShapes 不支持 add_shape，
    直接向 spTree 追加 p:sp 元素实现。
    """
    sp_tree = owner._element.spTree
    n = len(sp_tree.findall(qn("p:sp"))) + 1
    emu = lambda v: str(Inches(v))
    xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:nvSpPr>'
        f'<p:cNvPr id="{100 + n}" name="{name or "bar_" + str(n)}"/>'
        '<p:cNvSpPr/><p:nvPr/>'
        '</p:nvSpPr>'
        '<p:spPr>'
        f'<a:xfrm><a:off x="{emu(left)}" y="{emu(top)}"/>'
        f'<a:ext cx="{emu(width)}" cy="{emu(height)}"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        '<a:ln><a:noFill/></a:ln>'
        '</p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        '</p:sp>'
    )
    sp = parse_xml(xml)
    sp_tree.append(sp)
    return sp


def _add_frame(owner, left, top, width, height, color, name):
    """透明填充的边框矩形（图表占位区域指示）。"""
    sp_tree = owner._element.spTree
    n = len(sp_tree.findall(qn("p:sp"))) + 1
    emu = lambda v: str(Inches(v))
    xml = (
        '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:nvSpPr>'
        f'<p:cNvPr id="{100 + n}" name="{name}"/>'
        '<p:cNvSpPr/><p:nvPr/>'
        '</p:nvSpPr>'
        '<p:spPr>'
        f'<a:xfrm><a:off x="{emu(left)}" y="{emu(top)}"/>'
        f'<a:ext cx="{emu(width)}" cy="{emu(height)}"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '<a:noFill/>'
        f'<a:ln w="19050"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln>'
        '</p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        '</p:sp>'
    )
    sp = parse_xml(xml)
    sp_tree.append(sp)
    return sp


def make_pptx_template(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    master = prs.slide_masters[0]
    # 母版深色背景
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = BG_DARK
    # 母版底部品牌条（所有页面可见）
    _add_bar(master, 0, 7.32, 13.333, 0.18, BLUE)

    layouts = prs.slide_layouts

    # ─── layout[0] 封面 ─────────────────────────────────────────
    cover = layouts[0]
    cover.background.fill.solid()
    cover.background.fill.fore_color.rgb = BG_DARK
    title_ph = cover.placeholders[0]
    _move_ph(title_ph, 1.4, 2.0, 10.5, 1.8)
    _style_placeholder(title_ph, 44, WHITE, True, PP_ALIGN.LEFT)
    sub_ph = cover.placeholders[1]
    _move_ph(sub_ph, 1.5, 4.1, 8.0, 0.8)
    _style_placeholder(sub_ph, 20, CYAN, False, PP_ALIGN.LEFT)
    # 左侧发光条 + 装饰线
    _add_bar(cover, 0, 0, 0.08, 7.5, BLUE)
    _add_bar(cover, 1.5, 3.75, 3.0, 0.04, CYAN)
    _add_bar(cover, 1.5, 5.9, 10.3, 0.02, BLUE)

    # ─── layout[1] 内容页（目录/要点/表格/KPI 共用）─────────────
    content = layouts[1]
    title_ph = content.placeholders[0]
    _move_ph(title_ph, 0.8, 0.35, 11.7, 0.8)
    _style_placeholder(title_ph, 28, WHITE, True, PP_ALIGN.LEFT)
    body_ph = content.placeholders[1]
    _move_ph(body_ph, 0.9, 1.7, 11.5, 5.4)
    _style_placeholder(body_ph, 16, GRAY, False, PP_ALIGN.LEFT)
    _add_bar(content, 0.8, 1.25, 2.0, 0.045, CYAN)

    # ─── layout[2] 章节分隔页 ───────────────────────────────────
    section = layouts[2]
    title_ph = section.placeholders[0]
    _move_ph(title_ph, 1.5, 2.6, 10.3, 1.2)
    _style_placeholder(title_ph, 36, WHITE, True, PP_ALIGN.CENTER)
    body_ph = section.placeholders[1]
    _move_ph(body_ph, 1.5, 4.15, 10.3, 0.9)
    _style_placeholder(body_ph, 16, GRAY, False, PP_ALIGN.CENTER)
    _add_bar(section, 5.5, 3.95, 2.3, 0.045, CYAN)

    # ─── layout[5] 图表页（仅标题，图表图片运行时插入）─────────
    chart = layouts[5]
    title_ph = chart.placeholders[0]
    _move_ph(title_ph, 0.8, 0.35, 11.7, 0.8)
    _style_placeholder(title_ph, 26, WHITE, True, PP_ALIGN.LEFT)
    _add_bar(chart, 0.8, 1.25, 2.0, 0.045, CYAN)
    # 图表占位区域边框（指示图表区几何，运行时按此区域插入图片）
    _add_frame(chart, 1.0, 1.5, 11.3, 5.6, BG_CARD, "chart_frame")

    # ─── layout[6] 致谢（空白版式，仅深色背景）──────────────────
    thanks = layouts[6]
    thanks.background.fill.solid()
    thanks.background.fill.fore_color.rgb = BG_DARK

    prs.save(path)
    print(f"[build_templates] PPT 模板已生成: {path}")


# ═══════════════════════════════════════════════════════════════
# Word 模板
# ═══════════════════════════════════════════════════════════════

DOCX_HEADING1 = DRGBColor(0x1A, 0x23, 0x7E)   # 深蓝
DOCX_HEADING2 = DRGBColor(0x28, 0x35, 0x93)
DOCX_PRIMARY = DRGBColor(0x00, 0x71, 0xB8)    # 科技蓝
DOCX_BODY = DRGBColor(0x2C, 0x2C, 0x2C)
DOCX_GRAY = DRGBColor(0x80, 0x80, 0x80)


def _docx_set_style(doc, name, font_name, size, color, bold=False):
    style = doc.styles[name]
    style.font.name = font_name
    style.font.size = DPt(size)
    style.font.color.rgb = color
    style.font.bold = bold
    if style.element.rPr is not None and style.element.rPr.rFonts is not None:
        style.element.rPr.rFonts.set(dqn("w:eastAsia"), font_name)


def _setup_docx_styles(doc):
    _docx_set_style(doc, "Normal", EA_FONT, 11, DOCX_BODY)
    _docx_set_style(doc, "Heading 1", EA_FONT, 22, DOCX_HEADING1, True)
    _docx_set_style(doc, "Heading 2", EA_FONT, 16, DOCX_HEADING2, True)
    _docx_set_style(doc, "Heading 3", EA_FONT, 13, DOCX_PRIMARY, True)
    _docx_set_style(doc, "Heading 4", EA_FONT, 12, DOCX_PRIMARY, True)
    _docx_set_style(doc, "List Bullet", EA_FONT, 11, DOCX_BODY)
    _docx_set_style(doc, "List Number", EA_FONT, 11, DOCX_BODY)
    for sec in doc.sections:
        sec.top_margin = docx_cm(2.5)
        sec.bottom_margin = docx_cm(2.5)
        sec.left_margin = docx_cm(2.8)
        sec.right_margin = docx_cm(2.8)


def docx_cm(v):
    from docx.shared import Cm
    return Cm(v)


def _add_toc_field(paragraph):
    """在段落中插入 TOC 域（打开 Word 后 F9 更新生成目录）。"""
    run = paragraph.add_run()
    fld_char = OxmlElement("w:fldChar")
    fld_char.set(dqn("w:fldCharType"), "begin")
    run._r.append(fld_char)
    run2 = paragraph.add_run()
    instr = OxmlElement("w:instrText")
    instr.set(dqn("xml:space"), "preserve")
    instr.text = ' TOC \\o "1-3" \\h \\z \\u '
    run2._r.append(instr)
    run3 = paragraph.add_run()
    fld_char2 = OxmlElement("w:fldChar")
    fld_char2.set(dqn("w:fldCharType"), "end")
    run3._r.append(fld_char2)


def _docx_para(doc, text="", style=None, size=None, color=None, bold=False,
               align=None, italic=False, font_name=EA_FONT):
    p = doc.add_paragraph(style=style)
    if align is not None:
        p.alignment = align
    run = p.add_run(text)
    if size:
        run.font.size = DPt(size)
    if color:
        run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run._element.rPr.rFonts.set(dqn("w:eastAsia"), font_name)
    return p


def _setup_docx_header_footer(doc):
    section = doc.sections[0]
    header = section.header
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = hp.add_run("数据分析报告")
    run.font.size = DPt(9)
    run.font.color.rgb = DOCX_GRAY
    run.font.name = EA_FONT
    run._element.rPr.rFonts.set(dqn("w:eastAsia"), EA_FONT)

    footer = section.footer
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = fp.add_run()
    fld = OxmlElement("w:fldSimple")
    fld.set(dqn("w:instr"), "PAGE")
    run._r.append(fld)
    run.font.size = DPt(9)
    run.font.color.rgb = DOCX_GRAY


def make_docx_template(path):
    doc = Document()
    _setup_docx_styles(doc)

    # ─── 封面 ──────────────────────────────────────────────────
    for _ in range(4):
        doc.add_paragraph()
    p = _docx_para(doc, "数据分析报告", size=12, color=DOCX_PRIMARY,
                   align=WD_ALIGN_PARAGRAPH.CENTER, bold=True)
    p = _docx_para(doc, "{{ title }}", size=28, color=DOCX_HEADING1,
                   bold=True, align=WD_ALIGN_PARAGRAPH.CENTER)
    _docx_para(doc, "{{ subtitle }}", size=14, color=DOCX_PRIMARY,
               align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    doc.add_paragraph()
    for line in ("{{ org }}", "{{ author }}", "{{ date }}"):
        _docx_para(doc, line, size=12, color=DOCX_GRAY,
                   align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_page_break()

    # ─── 目录页 ────────────────────────────────────────────────
    _docx_para(doc, "目  录", size=18, color=DOCX_HEADING1, bold=True,
               align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    _add_toc_field(doc.add_paragraph())
    _docx_para(doc, "（在 Word 中右键此处 → 更新域，以生成目录）",
               size=9, color=DOCX_GRAY, italic=True,
               align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_page_break()

    # ─── 正文骨架（docxtpl 标签）───────────────────────────
    # 注意：docxtpl 0.20.x 中 {%p %} 段落级标签必须独占段落——
    # 含标签的段落会被整体替换为裸 Jinja 标签，段内其他内容丢失。
    # 因此 if/for/endif/endfor 标签各自成段，被条件/循环的正文单独成段。
    # Tiny 段（1pt）为标签段，渲染后由 word_generator 后处理清理。
    skeleton = [
        ("{% for section in sections %}", "Tiny"),
        ("{%p if section.level == 1 %}", "Tiny"),
        ("{{ section.title }}", "Heading 1"),
        ("{%p endif %}", "Tiny"),
        ("{%p if section.level == 2 %}", "Tiny"),
        ("{{ section.title }}", "Heading 2"),
        ("{%p endif %}", "Tiny"),
        ("{%p if section.level == 3 %}", "Tiny"),
        ("{{ section.title }}", "Heading 3"),
        ("{%p endif %}", "Tiny"),
        ("{% for block in section.blocks %}", "Tiny"),
        ("{%p if block.kind == 'paragraph' %}", "Tiny"),
        ("{{ block.content }}", "Normal"),
        ("{%p endif %}", "Tiny"),
        ("{%p if block.kind == 'subheading' %}", "Tiny"),
        ("{{ block.content }}", "Heading 4"),
        ("{%p endif %}", "Tiny"),
        ("{%p for item in block['items'] %}", "Tiny"),
        ("{{ item }}", "List Bullet"),
        ("{%p endfor %}", "Tiny"),
        ("{%p for line in block['lines'] %}", "Tiny"),
        ("{{ line }}", "CodeBlock"),
        ("{%p endfor %}", "Tiny"),
        ("{%p if block.kind == 'table' %}", "Tiny"),
        ("%%WEBSQL_TABLE:{{ block.index }}%%", "Normal"),
        ("{%p endif %}", "Tiny"),
        ("{%p if block.kind == 'image' %}", "Tiny"),
        ("%%WEBSQL_IMAGE:{{ block.index }}%%", "Normal"),
        ("{%p endif %}", "Tiny"),
        ("{% endfor %}", "Tiny"),
        ("{% endfor %}", "Tiny"),
    ]
    for text, style in skeleton:
        if style == "Tiny":
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.font.size = DPt(1)
        elif style == "CodeBlock":
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.font.size = DPt(9)
            run.font.name = "Courier New"
            run._element.rPr.rFonts.set(dqn("w:eastAsia"), EA_FONT)
            pPr = p._element.get_or_add_pPr()
            shd = pPr.makeelement(dqn("w:shd"), {
                dqn("w:val"): "clear", dqn("w:color"): "auto",
                dqn("w:fill"): "F5F5F5"})
            pPr.append(shd)
        else:
            doc.add_paragraph(text, style=style)

    _setup_docx_header_footer(doc)
    doc.save(path)
    print(f"[build_templates] Word 模板已生成: {path}")


if __name__ == "__main__":
    base = os.path.dirname(os.path.abspath(__file__))
    ppt_dir = os.path.join(base, "..", "export-ppt", "templates")
    word_dir = os.path.join(base, "..", "export-word", "templates")
    os.makedirs(ppt_dir, exist_ok=True)
    os.makedirs(word_dir, exist_ok=True)
    make_pptx_template(os.path.join(ppt_dir, "slides_template.pptx"))
    make_docx_template(os.path.join(word_dir, "report_template.docx"))
